import streamlit as st
from deep_translator import GoogleTranslator
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PyPDF2 import PdfReader
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
import tempfile
import os

# Supported Languages for Deep Translator
LANGUAGES = [
    'af', 'sq', 'ar', 'bn', 'bg', 'zh-CN', 'zh-TW', 'hr', 'cs', 'da', 'nl', 'en', 'et',
    'tl', 'fi', 'fr', 'de', 'el', 'gu', 'he', 'hi', 'hu', 'id', 'it', 'ja', 'kn', 'ko',
    'lv', 'lt', 'ms', 'ml', 'mr', 'ne', 'no', 'fa', 'pl', 'pt', 'pa', 'ro', 'ru', 'sr',
    'si', 'sk', 'sl', 'es', 'sv', 'ta', 'te', 'th', 'tr', 'uk', 'ur', 'vi'
]

LANGUAGE_NAMES = {
    'af': 'Afrikaans', 'sq': 'Albanian', 'ar': 'Arabic', 'bn': 'Bengali', 'bg': 'Bulgarian',
    'zh-CN': 'Chinese (Simplified)', 'zh-TW': 'Chinese (Traditional)', 'hr': 'Croatian',
    'cs': 'Czech', 'da': 'Danish', 'nl': 'Dutch', 'en': 'English', 'et': 'Estonian', 'tl': 'Filipino',
    'fi': 'Finnish', 'fr': 'French', 'de': 'German', 'el': 'Greek', 'gu': 'Gujarati', 'he': 'Hebrew',
    'hi': 'Hindi', 'hu': 'Hungarian', 'id': 'Indonesian', 'it': 'Italian', 'ja': 'Japanese',
    'kn': 'Kannada', 'ko': 'Korean', 'lv': 'Latvian', 'lt': 'Lithuanian', 'ms': 'Malay',
    'ml': 'Malayalam', 'mr': 'Marathi', 'ne': 'Nepali', 'no': 'Norwegian', 'fa': 'Persian',
    'pl': 'Polish', 'pt': 'Portuguese', 'pa': 'Punjabi', 'ro': 'Romanian', 'ru': 'Russian',
    'sr': 'Serbian', 'si': 'Sinhala', 'sk': 'Slovak', 'sl': 'Slovenian', 'es': 'Spanish',
    'sv': 'Swedish', 'ta': 'Tamil', 'te': 'Telugu', 'th': 'Thai', 'tr': 'Turkish', 'uk': 'Ukrainian',
    'ur': 'Urdu', 'vi': 'Vietnamese'
}

st.set_page_config(page_title="File Translator", layout="centered")

st.title(":earth_africa: Universal File Translator")
st.write("Upload your file, select input & output language, and instantly get a translated version!")

# Upload section
uploaded_file = st.file_uploader("Upload a file (Word, PowerPoint, Excel, or PDF)", type=["docx", "pptx", "xlsx", "pdf"])

# Language selectors
col1, col2 = st.columns(2)
with col1:
    src_lang_name = st.selectbox("Select Input Language", options=[LANGUAGE_NAMES[code] for code in LANGUAGES], index=[LANGUAGE_NAMES[code] for code in LANGUAGES].index('English'))
with col2:
    tgt_lang_name = st.selectbox("Select Output Language", options=[LANGUAGE_NAMES[code] for code in LANGUAGES], index=[LANGUAGE_NAMES[code] for code in LANGUAGES].index('French'))

src_lang = [code for code, name in LANGUAGE_NAMES.items() if name == src_lang_name][0]
tgt_lang = [code for code, name in LANGUAGE_NAMES.items() if name == tgt_lang_name][0]

if uploaded_file:
    st.info(f"File uploaded: {uploaded_file.name}")

    if st.button(":arrows_counterclockwise: Translate File"):
        with st.spinner("Translating file... Please wait :hourglass_flowing_sand:"):
            with tempfile.NamedTemporaryFile(delete=False) as tmp:
                tmp.write(uploaded_file.read())
                tmp_path = tmp.name

            ext = os.path.splitext(uploaded_file.name)[1].lower()
            translated_path = f"translated_{uploaded_file.name}"

            try:
                translator = GoogleTranslator(source=src_lang, target=tgt_lang)

                if ext == ".docx":
                    doc = Document(tmp_path)
                    for para in doc.paragraphs:
                        if para.text.strip():
                            para.text = translator.translate(para.text)
                    doc.save(translated_path)

                elif ext == ".pptx":
                    prs = Presentation(tmp_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                shape.text = translator.translate(shape.text)
                    prs.save(translated_path)

                elif ext == ".xlsx":
                    wb = load_workbook(tmp_path)
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        for row in ws.iter_rows():
                            for cell in row:
                                if cell.value and isinstance(cell.value, str):
                                    cell.value = translator.translate(cell.value)
                    wb.save(translated_path)

                elif ext == ".pdf":
                    reader = PdfReader(tmp_path)
                    c = canvas.Canvas(translated_path, pagesize=A4)
                    width, height = A4
                    for page in reader.pages:
                        text = page.extract_text()
                        if text:
                            translated_text = translator.translate(text)
                            y = height - 50
                            for line in translated_text.split("\n"):
                                c.drawString(40, y, line[:1000])
                                y -= 15
                                if y < 50:
                                    c.showPage()
                                    y = height - 50
                            c.showPage()
                    c.save()

                else:
                    st.error("Unsupported file type.")

                # Download translated file
                with open(translated_path, "rb") as f:
                    st.success(":white_check_mark: Translation completed successfully!")
                    st.download_button(
                        ":inbox_tray: Download Translated File",
                        data=f,
                        file_name=translated_path,
                        mime="application/octet-stream"
                    )

            except Exception as e:
                st.error(f"Error: {e}")

            finally:
                os.remove(tmp_path)
